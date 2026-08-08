#!/usr/bin/env python3
"""Build clean learner visuals and source-task manifests from New HSK PPTX.

Only useful PPT illustrations are exported. PDF pages remain trace-only metadata
and are never copied into the learner-facing asset folder.
"""
from __future__ import annotations

import argparse
import io
import json
import re
import shutil
from pathlib import Path
from typing import Any, Iterable

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from course_config import lesson_count, lesson_numbers

EMU_PER_INCH = 914400


def read_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def slide_text(slide) -> str:
    return " | ".join(
        " ".join(getattr(shape, "text", "").split())
        for shape in slide.shapes
        if getattr(shape, "text", "").strip()
    )


def shape_box(shape) -> tuple[float, float, float, float]:
    return (
        shape.left / EMU_PER_INCH,
        shape.top / EMU_PER_INCH,
        shape.width / EMU_PER_INCH,
        shape.height / EMU_PER_INCH,
    )


def find_target_slides(prs: Presentation) -> dict[str, int]:
    targets: dict[str, int] = {}
    for index, slide in enumerate(prs.slides, 1):
        text = slide_text(slide)
        compact = text.replace(" ", "")
        if "warmup" not in targets and index <= 8:
            lower_text = text.lower()
            has_warmup_marker = "warm-up" in lower_text or "warmup" in lower_text or "热身" in compact
            has_first_task = any(token in compact for token in ("1.", "1。", "1．"))
            if has_warmup_marker and has_first_task:
                targets["warmup"] = index
        for number in (1, 2, 3, 4):
            key = f"text{number}"
            if key not in targets and f"课文{number}" in compact and ("课文内容请见" in compact or "Pleaserefertopage" in compact):
                targets[key] = index
        if "activity" not in targets and "课堂活动" in compact:
            targets["activity"] = index
        if "summary" not in targets and "学习小结" in compact:
            targets["summary"] = index
        if "extension" not in targets and "小语的彩蛋" in compact:
            targets["extension"] = index
    return targets


def standalone_tokens(slide, pattern: str) -> list[tuple[str, float, float, float, float]]:
    regex = re.compile(pattern)
    rows: list[tuple[str, float, float, float, float]] = []
    for shape in slide.shapes:
        text = " ".join(getattr(shape, "text", "").split()).strip()
        if regex.fullmatch(text):
            x, y, w, h = shape_box(shape)
            rows.append((text, x, y, w, h))
    rows.sort(key=lambda row: (round(row[2], 1), row[1]))
    return rows


def picture_shapes(shapes: Iterable[Any]) -> list[Any]:
    direct: list[Any] = []
    groups: list[Any] = []
    for shape in shapes:
        # PowerPoint picture placeholders are exposed by python-pptx as
        # PLACEHOLDER instead of PICTURE, but still carry an .image payload.
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
            direct.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            groups.append(shape)
    direct.sort(key=lambda shape: (shape.top, shape.left))
    result = list(direct)
    for group in sorted(groups, key=lambda shape: (shape.top, shape.left)):
        children = picture_shapes(group.shapes)
        children.sort(key=lambda shape: (shape.top, shape.left))
        result.extend(children)
    return result


def image_from_picture(shape) -> Image.Image:
    image = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
    left = max(0.0, float(shape.crop_left or 0.0))
    top = max(0.0, float(shape.crop_top or 0.0))
    right = max(0.0, float(shape.crop_right or 0.0))
    bottom = max(0.0, float(shape.crop_bottom or 0.0))
    x1 = int(round(left * image.width))
    y1 = int(round(top * image.height))
    x2 = int(round((1.0 - right) * image.width))
    y2 = int(round((1.0 - bottom) * image.height))
    if x2 > x1 and y2 > y1:
        image = image.crop((x1, y1, x2, y2))
    return image


def trim_white(image: Image.Image, tolerance: int = 247, padding: int = 12) -> Image.Image:
    gray = image.convert("L")
    inverted = ImageOps.invert(gray)
    mask = inverted.point(lambda value: 255 if value > (255 - tolerance) else 0)
    box = mask.getbbox()
    if not box:
        return image
    x1, y1, x2, y2 = box
    return image.crop((max(0, x1 - padding), max(0, y1 - padding), min(image.width, x2 + padding), min(image.height, y2 + padding)))


def contain(image: Image.Image, width: int, height: int) -> Image.Image:
    canvas = Image.new("RGB", (width, height), "white")
    clone = image.copy()
    clone.thumbnail((width - 24, height - 24), Image.Resampling.LANCZOS)
    x = (width - clone.width) // 2
    y = (height - clone.height) // 2
    canvas.paste(clone, (x, y))
    return canvas


def warmup_visual(slide) -> tuple[Image.Image | None, list[str]]:
    answer_tokens = standalone_tokens(slide, r"[A-F]")
    answers = [row[0] for row in answer_tokens]
    pictures = picture_shapes(slide.shapes)
    if not pictures:
        return None, answers
    images = [trim_white(image_from_picture(shape), padding=6) for shape in pictures]
    expected = len(answers)
    if expected and len(images) > expected:
        # Decorative marks are usually tiny; keep the largest expected images.
        images = sorted(images, key=lambda img: img.width * img.height, reverse=True)[:expected]
    if len(images) == 1:
        return trim_white(images[0]), answers
    count = len(images)
    columns = 3 if count == 6 else min(4, count)
    rows = (count + columns - 1) // columns
    card_w, card_h, gap = 300, 230, 18
    canvas = Image.new("RGB", (columns * card_w + (columns - 1) * gap, rows * card_h + (rows - 1) * gap), "white")
    for index, image in enumerate(images):
        card = contain(image, card_w, card_h)
        x = (index % columns) * (card_w + gap)
        y = (index // columns) * (card_h + gap)
        canvas.paste(card, (x, y))
    return trim_white(canvas, padding=4), answers


def largest_picture(slide, min_area: float = 1.0):
    candidates = []
    for shape in picture_shapes(slide.shapes):
        x, y, w, h = shape_box(shape)
        area = w * h
        if area >= min_area and y >= 1.1:
            candidates.append((area, shape))
    if not candidates:
        return None
    return max(candidates, key=lambda item: item[0])[1]


def save_webp(image: Image.Image, path: Path, max_width: int = 1400) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    image = trim_white(image.convert("RGB"), padding=8)
    if image.width > max_width:
        height = int(round(image.height * max_width / image.width))
        image = image.resize((max_width, height), Image.Resampling.LANCZOS)
    image.save(path, "WEBP", quality=88, method=6)


def is_image_match_warmup(text: str) -> bool:
    compact = re.sub(r"[\s|]+", "", text).lower()
    return (
        ("选择对应的图片" in compact or "选择相应的图片" in compact)
        or ("matchthewords" in compact and "pictures" in compact)
    )


def listening_answers(slide) -> tuple[str, list[str]] | None:
    tokens = standalone_tokens(slide, r"[ABC√×]")
    if not tokens:
        return None
    answers = [token[0] for token in tokens]
    if all(value in {"√", "×"} for value in answers):
        return "listening-tf", answers
    return "listening-mcq", answers


def lesson_section_map(lesson: dict[str, Any]) -> dict[str, dict[str, Any]]:
    result: dict[str, dict[str, Any]] = {}
    for section in lesson.get("entities", {}).get("contentSections", []):
        title = section.get("title", "")
        kind = section.get("kind", "")
        if kind == "warmup":
            result["warmup"] = section
        elif kind == "activity":
            result["activity"] = section
        elif kind == "passage":
            # HSK 2-3 use the fourth textbook track (x-7) as the passage/text 4.
            result["text4"] = section
        elif title.lower().startswith("tổng kết"):
            result["summary"] = section
        elif kind == "extension":
            result["extension"] = section
        else:
            match = re.search(r"Bài khóa\s*(\d+)", title, re.I)
            if match:
                result[f"text{match.group(1)}"] = section
    return result


def page_for_section(lesson: dict[str, Any], section: dict[str, Any]) -> int | None:
    traces = lesson.get("source", {}).get("pageTrace", []) or []
    title = section.get("title", "")
    kind = section.get("kind", "")
    needles: list[str] = []
    match = re.search(r"Bài khóa\s*(\d+)", title, re.I)
    if match:
        needles = [f"Bài khóa {match.group(1)}"]
    elif kind in {"warmup", "objectives"}:
        needles = ["khởi động", "Tên bài"]
    elif kind == "exercise":
        needles = ["Bài tập tổng hợp", "Bài tập"]
    elif kind == "activity":
        needles = ["Hoạt động"]
    elif kind == "extension":
        needles = ["Món quà", "mở rộng"]
    elif title.lower().startswith("tổng kết"):
        needles = ["Tổng kết", "tự đánh giá"]
    for needle in needles:
        for trace in traces:
            if needle.lower() in str(trace.get("content", "")).lower():
                return int(trace.get("pdfPage"))
    if traces:
        return int(traces[0 if kind in {"warmup", "objectives"} else -1].get("pdfPage"))
    return None


def add_pdf_trace(rows: list[dict[str, Any]], lesson: dict[str, Any], section: dict[str, Any], level: int) -> None:
    page = page_for_section(lesson, section)
    if page is None:
        return
    rows.append({
        "id": f"{lesson['id']}-visual-pdf-{page}",
        "alt": f"Trang nguồn Bài {lesson['lessonNumber']}",
        "caption": f"Trang nguồn PDF {page}",
        "sourceType": "pdf",
        "sourceRef": f"新HSK{level} 教材 - trang PDF {page}",
        "assetPolicy": "trace-only",
    })


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--repo", type=Path, required=True)
    parser.add_argument("--level", type=int, required=True)
    parser.add_argument("--ppt-dir", type=Path, required=True)
    parser.add_argument("--lesson-start", type=int, default=1)
    parser.add_argument("--lesson-end", type=int, default=None)
    args = parser.parse_args()

    repo = args.repo.resolve()
    level = args.level
    data_dir = repo / f"modules/new-hsk-course/data/hsk{level}"
    source_dir = repo / f"modules/new-hsk-course/source/hsk{level}"
    asset_root = repo / f"modules/new-hsk-course/assets/visuals/hsk{level}"
    source_dir.mkdir(parents=True, exist_ok=True)
    asset_root.mkdir(parents=True, exist_ok=True)

    ppt_by_lesson: dict[int, Path] = {}
    for path in args.ppt_dir.glob("*.pptx"):
        match = re.search(r"第\s*(\d+)\s*课", path.name)
        if match:
            ppt_by_lesson[int(match.group(1))] = path

    visual_manifest: dict[str, Any] = {"version": 1, "course": "new-hsk-course", "level": level, "lessons": {}}
    task_manifest: dict[str, Any] = {"version": 1, "course": "new-hsk-course", "level": level, "lessons": {}}

    lesson_end = args.lesson_end if args.lesson_end is not None else lesson_count(level)
    selected_lessons = range(max(1, args.lesson_start), min(lesson_count(level), lesson_end) + 1)

    for lesson_number in selected_lessons:
        pptx_path = ppt_by_lesson.get(lesson_number)
        if not pptx_path:
            raise FileNotFoundError(f"Missing PPT for lesson {lesson_number}")
        lesson = read_json(data_dir / f"lesson-{lesson_number:02d}.json")
        section_map = lesson_section_map(lesson)
        prs = Presentation(pptx_path)
        targets = find_target_slides(prs)
        lesson_assets = asset_root / f"lesson-{lesson_number:02d}"
        if lesson_assets.exists():
            shutil.rmtree(lesson_assets)
        lesson_assets.mkdir(parents=True)

        visual_sections: dict[str, list[dict[str, Any]]] = {}
        task_sections: dict[str, list[dict[str, Any]]] = {}
        for section in lesson.get("entities", {}).get("contentSections", []):
            rows: list[dict[str, Any]] = []
            add_pdf_trace(rows, lesson, section, level)
            if rows:
                visual_sections[section["id"]] = rows

        warmup_section = section_map.get("warmup")
        warmup_slide = targets.get("warmup")
        if warmup_section and warmup_slide:
            image, answers = warmup_visual(prs.slides[warmup_slide - 1])
            if image:
                rel = f"assets/visuals/hsk{level}/lesson-{lesson_number:02d}/ppt-warmup.webp"
                save_webp(image, repo / "modules/new-hsk-course" / rel)
                visual_sections.setdefault(warmup_section["id"], []).insert(0, {
                    "id": f"{lesson['id']}-visual-ppt-warmup",
                    "src": rel,
                    "alt": f"Hình khởi động Bài {lesson_number}",
                    "caption": "Hình khởi động",
                    "sourceType": "ppt",
                    "sourceRef": f"PPT Bài {lesson_number} - slide {warmup_slide}",
                    "assetPolicy": "learner-visual",
                })
            warmup_text = slide_text(prs.slides[warmup_slide - 1])
            if answers and is_image_match_warmup(warmup_text):
                task_sections[warmup_section["id"]] = [{
                    "id": f"{lesson['id']}-task-warmup",
                    "type": "image-match",
                    "answers": answers,
                    "sourceRef": f"PPT Bài {lesson_number} - Khởi động slide {warmup_slide}",
                }]

        for text_number in (1, 2, 3, 4):
            key = f"text{text_number}"
            section = section_map.get(key)
            slide_number = targets.get(key)
            if not section or not slide_number:
                continue
            shape = largest_picture(prs.slides[slide_number - 1], min_area=1.5)
            if shape:
                rel = f"assets/visuals/hsk{level}/lesson-{lesson_number:02d}/ppt-text-{text_number}.webp"
                save_webp(image_from_picture(shape), repo / "modules/new-hsk-course" / rel)
                visual_sections.setdefault(section["id"], []).insert(0, {
                    "id": f"{lesson['id']}-visual-ppt-text-{text_number}",
                    "src": rel,
                    "alt": f"Hình tình huống Bài khóa {text_number} - Bài {lesson_number}",
                    "caption": f"Hình tình huống Bài khóa {text_number}",
                    "sourceType": "ppt",
                    "sourceRef": f"PPT Bài {lesson_number} - slide {slide_number}",
                    "assetPolicy": "learner-visual",
                })

        activity_section = section_map.get("activity")
        activity_slide = targets.get("activity")
        if activity_section and activity_slide:
            shape = largest_picture(prs.slides[activity_slide - 1], min_area=1.0)
            if shape:
                rel = f"assets/visuals/hsk{level}/lesson-{lesson_number:02d}/ppt-activity.webp"
                save_webp(image_from_picture(shape), repo / "modules/new-hsk-course" / rel)
                visual_sections.setdefault(activity_section["id"], []).insert(0, {
                    "id": f"{lesson['id']}-visual-ppt-activity",
                    "src": rel,
                    "alt": f"Hình hoạt động trên lớp - Bài {lesson_number}",
                    "caption": "Hình hoạt động trên lớp",
                    "sourceType": "ppt",
                    "sourceRef": f"PPT Bài {lesson_number} - slide {activity_slide}",
                    "assetPolicy": "learner-visual",
                })

        for text_number in (1, 2, 3, 4):
            key = f"text{text_number}"
            section = section_map.get(key)
            intro_slide = targets.get(key)
            if not section or not intro_slide or intro_slide >= len(prs.slides):
                continue
            question_slide = intro_slide + 1
            parsed = listening_answers(prs.slides[question_slide - 1])
            if not parsed:
                continue
            task_type, answers = parsed
            task_sections.setdefault(section["id"], []).append({
                "id": f"{lesson['id']}-task-listen-{text_number:02d}",
                "type": task_type,
                "answers": answers,
                "audioRef": f"{lesson_number}-{text_number * 2 - 1}",
                "sourceRef": f"PPT Bài {lesson_number} - đáp án nghe slide {question_slide}",
            })

        visual_manifest["lessons"][lesson["id"]] = {"sections": visual_sections}
        task_manifest["lessons"][lesson["id"]] = {"sections": task_sections}
        print(f"lesson {lesson_number:02d}: visuals={sum(map(len, visual_sections.values()))}, tasks={sum(map(len, task_sections.values()))}, slides={targets}")

    (source_dir / "visual-manifest.json").write_text(json.dumps(visual_manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    (source_dir / "source-task-manifest.json").write_text(json.dumps(task_manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print("wrote", source_dir / "visual-manifest.json")
    print("wrote", source_dir / "source-task-manifest.json")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
